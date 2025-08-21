import os
import logging
import fitz  # PyMuPDF
import docx
import pandas as pd  # для Excel и CSV
from pptx import Presentation  # для PPT/PPTX


logger = logging.getLogger(__name__)


def read_text_file(file_path: str) -> str | None:
    """Прочитать содержимое текстового файла (ограничение по длине для скорости)."""
    max_chars = 3000
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            return file.read(max_chars)
    except Exception as e:
        logger.error("Error reading text file %s: %s", file_path, e)
        return None


def read_docx_file(file_path: str) -> str | None:
    """Прочитать текст из .docx/.doc (через python-docx)."""
    try:
        doc = docx.Document(file_path)
        return '\n'.join(p.text for p in doc.paragraphs)
    except Exception as e:
        logger.error("Error reading DOCX file %s: %s", file_path, e)
        return None


def read_pdf_file(file_path: str) -> str | None:
    """Прочитать текст из первых страниц PDF (через PyMuPDF)."""
    try:
        doc = fitz.open(file_path)
        num_pages_to_read = 3  # можно настроить
        pages = []
        for i in range(min(num_pages_to_read, len(doc))):
            page = doc.load_page(i)
            pages.append(page.get_text())
        return '\n'.join(pages)
    except Exception as e:
        logger.error("Error reading PDF file %s: %s", file_path, e)
        return None


def read_spreadsheet_file(file_path: str) -> str | None:
    """Прочитать таблицу из Excel/CSV и вернуть как текст."""
    try:
        if file_path.lower().endswith('.csv'):
            df = pd.read_csv(file_path)
        else:
            df = pd.read_excel(file_path)
        return df.to_string()
    except Exception as e:
        logger.error("Error reading spreadsheet file %s: %s", file_path, e)
        return None


def read_ppt_file(file_path: str) -> str | None:
    """Прочитать текст из слайдов PowerPoint."""
    try:
        prs = Presentation(file_path)
        chunks = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    chunks.append(shape.text)
        return '\n'.join(chunks)
    except Exception as e:
        logger.error("Error reading PowerPoint file %s: %s", file_path, e)
        return None


def read_file_data(file_path: str) -> str | None:
    """Выбрать способ чтения по расширению файла и вернуть текст."""
    ext = os.path.splitext(file_path.lower())[1]
    if ext in ('.txt', '.md'):
        return read_text_file(file_path)
    if ext in ('.docx', '.doc'):
        return read_docx_file(file_path)
    if ext == '.pdf':
        return read_pdf_file(file_path)
    if ext in ('.xls', '.xlsx', '.csv'):
        return read_spreadsheet_file(file_path)
    if ext in ('.ppt', '.pptx'):
        return read_ppt_file(file_path)
    return None  # Неподдерживаемый тип


def collect_file_paths(base_path: str) -> list[str]:
    """Собрать все пути файлов из директории (или вернуть один путь), игнорируя скрытые файлы."""
    if os.path.isfile(base_path):
        return [base_path]
    file_paths: list[str] = []
    for root, _, files in os.walk(base_path):
        for name in files:
            if not name.startswith('.'):  # исключаем скрытые
                file_paths.append(os.path.join(root, name))
    return file_paths
